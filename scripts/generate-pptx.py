#!/usr/bin/env python3
"""
generate-pptx.py -- Generate a PowerPoint file with all OCHA humanitarian icons.

Icons are organized by family (as defined in metadata.json) and inserted as
SVG images with a PNG fallback.  Modern PowerPoint (Office 365 / 2019+) will
display the SVG natively; the user can right-click any icon and choose
"Convert to Shape" to get fully editable, color-changeable vector shapes.

Dependencies (all in .venv):
    python-pptx >= 1.0
    lxml
    Pillow
"""

import io
import json
import math
import os
import re
import struct
import sys
import zlib

from lxml import etree
from PIL import Image as PILImage, ImageDraw

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.opc.package import Part
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.packuri import PackURI

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------
ICON_COLOR = "#009edb"          # OCHA / UN blue
ICON_SIZE_IN = 0.55             # icon square size in inches
LABEL_FONT_SIZE = Pt(7)         # font size for icon label
HEADER_FONT_SIZE = Pt(14)       # font size for family header
COLS = 10                       # icons per row
ROW_HEIGHT_IN = 0.95            # height of one icon row (icon + label + gap)
COL_WIDTH_IN = 1.15             # width of one icon column
HEADER_HEIGHT_IN = 0.45         # family header height
PAGE_TOP_MARGIN_IN = 0.5        # top margin
PAGE_LEFT_MARGIN_IN = 0.45      # left margin
PAGE_BOTTOM_MARGIN_IN = 0.35    # bottom margin reserved
SLIDE_WIDTH_IN = 13.333         # widescreen
SLIDE_HEIGHT_IN = 7.5

# Derived
USABLE_HEIGHT_IN = SLIDE_HEIGHT_IN - PAGE_TOP_MARGIN_IN - PAGE_BOTTOM_MARGIN_IN
ICON_LABEL_COLOR = RGBColor(0x4A, 0x4A, 0x4A)
HEADER_TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)
HEADER_BG_COLOR = RGBColor(0xE8, 0xF4, 0xFB)  # light blue tint
UN_BLUE = RGBColor(0x00, 0x9E, 0xDB)

# Namespace constants for SVG-in-PPTX
ASVG_URI = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
REPO_ROOT = os.path.dirname(SCRIPT_DIR)
METADATA_PATH = os.path.join(REPO_ROOT, "metadata.json")
SVG_DIR = os.path.join(REPO_ROOT, "svg")
OUTPUT_DIR = os.path.join(REPO_ROOT, "output")
OUTPUT_PATH = os.path.join(OUTPUT_DIR, "Humanitarian_icons.pptx")

# Shape elements that should receive a fill
SHAPE_ELEMENTS = {"path", "circle", "rect", "polygon", "ellipse", "line", "polyline"}


# ---------------------------------------------------------------------------
# PNG fallback: create a small transparent PNG (1x1 pixel)
# ---------------------------------------------------------------------------
def make_fallback_png():
    """Create a minimal 1x1 transparent PNG as bytes."""
    def _chunk(chunk_type, data):
        c = chunk_type + data
        return (struct.pack(">I", len(data)) + c
                + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF))

    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 6, 0, 0, 0)
    raw = b"\x00" + bytes([0, 0, 0, 0])
    idat = zlib.compress(raw)

    png = b"\x89PNG\r\n\x1a\n"
    png += _chunk(b"IHDR", ihdr)
    png += _chunk(b"IDAT", idat)
    png += _chunk(b"IEND", b"")
    return png


# Pre-generate once
FALLBACK_PNG = make_fallback_png()


# ---------------------------------------------------------------------------
# SVG cleaning  (adapted from generate-grid.py)
# ---------------------------------------------------------------------------
def parse_viewbox(svg_text):
    """Return (min_x, min_y, width, height) from the viewBox attribute."""
    m = re.search(r'viewBox\s*=\s*"([^"]+)"', svg_text)
    if not m:
        return (0, 0, 48, 48)
    parts = m.group(1).split()
    return tuple(float(p) for p in parts)


def _extract_style_text(svg_text):
    """Extract the CSS <style> content (for evenodd detection)."""
    m = re.search(r"<style[^>]*>(.*?)</style>", svg_text, re.DOTALL)
    return m.group(1) if m else ""


def _extract_inner(svg_text):
    """Return the content between <svg ...> and </svg>, stripped of
    <?xml>, comments, <defs>, <style>, <title>."""
    t = re.sub(r"<\?xml[^?]*\?>", "", svg_text)
    t = re.sub(r"<!--.*?-->", "", t, flags=re.DOTALL)
    m = re.search(r"<svg[^>]*>(.*)</svg>", t, re.DOTALL)
    if not m:
        return ""
    inner = m.group(1)
    inner = re.sub(r"<defs>.*?</defs>", "", inner, flags=re.DOTALL)
    inner = re.sub(r"<defs\s[^>]*>.*?</defs>", "", inner, flags=re.DOTALL)
    inner = re.sub(r"<style[^>]*>.*?</style>", "", inner, flags=re.DOTALL)
    inner = re.sub(r"<title>.*?</title>", "", inner, flags=re.DOTALL)
    return inner.strip()


def _clean_inner(inner_svg, style_text, fill_color):
    """Resolve CSS classes to direct fill attributes."""
    # Determine which classes use fill-rule:evenodd
    evenodd_classes = set()
    if style_text:
        for m in re.finditer(
            r"\.([\w-]+)\s*\{[^}]*fill-rule\s*:\s*evenodd[^}]*\}", style_text
        ):
            evenodd_classes.add(m.group(1))

    def _process(match):
        tag = match.group(0)
        em = re.match(r"<(\w+)", tag)
        if not em:
            return tag
        elem = em.group(1)

        # Check evenodd before removing class
        needs_evenodd = False
        cls_m = re.search(r'class="([^"]*)"', tag)
        if cls_m and evenodd_classes:
            for c in cls_m.group(1).split():
                if c in evenodd_classes:
                    needs_evenodd = True
                    break

        # Strip class and style attributes
        tag = re.sub(r'\s+class="[^"]*"', "", tag)
        tag = re.sub(r'\s+style="[^"]*"', "", tag)

        if elem not in SHAPE_ELEMENTS:
            return tag

        if not re.search(r"\bfill\s*=", tag):
            tag = tag.replace(
                "<" + elem, '<' + elem + ' fill="' + fill_color + '"', 1
            )
        if needs_evenodd and "fill-rule" not in tag:
            tag = tag.replace(
                "<" + elem, "<" + elem + ' fill-rule="evenodd"', 1
            )
        return tag

    return re.sub(r"<(\w+)(?:\s[^>]*)?\s*/?>", _process, inner_svg)


def clean_svg(svg_path, fill_color=ICON_COLOR):
    """Read an SVG file and return a complete, self-contained, cleaned SVG
    string ready for embedding in PPTX."""
    with open(svg_path, "r", encoding="utf-8") as f:
        raw = f.read()

    _, _, vb_w, vb_h = parse_viewbox(raw)
    style = _extract_style_text(raw)
    inner = _extract_inner(raw)
    if not inner:
        return None, vb_w, vb_h

    cleaned = _clean_inner(inner, style, fill_color)

    # Rebuild a complete, self-contained SVG
    svg_doc = (
        f'<svg xmlns="http://www.w3.org/2000/svg" '
        f'viewBox="0 0 {vb_w} {vb_h}" '
        f'width="{vb_w}" height="{vb_h}">'
        f"{cleaned}</svg>"
    )
    return svg_doc, vb_w, vb_h


# ---------------------------------------------------------------------------
# Add an SVG picture to a slide  (PNG fallback + SVG via asvg:svgBlip)
# ---------------------------------------------------------------------------
_svg_counter = 0


def add_svg_picture(slide, svg_bytes, left, top, width, height, prs, desc=""):
    """Insert an SVG as a picture shape.

    Creates a PNG fallback image placeholder, then attaches the SVG data via
    the Office ``asvg:svgBlip`` extension so PowerPoint renders the SVG
    natively.
    """
    global _svg_counter
    _svg_counter += 1

    slide_part = slide.part

    # 1) Add picture using the tiny fallback PNG
    pic_shape = slide.shapes.add_picture(
        io.BytesIO(FALLBACK_PNG), left, top, width, height
    )

    # Update the description for accessibility
    if desc:
        pic_shape._element.find(qn("p:nvPicPr")).find(qn("p:cNvPr")).set(
            "descr", desc
        )

    # 2) Create an SVG part and relate it to the slide
    svg_part_name = PackURI(f"/ppt/media/image_svg{_svg_counter}.svg")
    svg_part = Part(
        svg_part_name,
        "image/svg+xml",
        blob=svg_bytes,
        package=prs.part.package,
    )
    svg_rId = slide_part.relate_to(svg_part, RT.IMAGE)

    # 3) Inject the <a:extLst><a:ext><asvg:svgBlip .../></a:ext></a:extLst>
    #    into the <a:blip> element of the picture
    blip = pic_shape._element.find(".//" + qn("a:blip"))
    ext_lst = etree.SubElement(blip, qn("a:extLst"))
    ext_el = etree.SubElement(ext_lst, qn("a:ext"))
    ext_el.set("uri", SVG_EXT_URI)
    svg_blip = etree.SubElement(ext_el, f"{{{ASVG_URI}}}svgBlip")
    svg_blip.set(qn("r:embed"), svg_rId)

    return pic_shape


# ---------------------------------------------------------------------------
# Add a text box helper
# ---------------------------------------------------------------------------
def add_textbox(slide, text, left, top, width, height, font_size=Pt(8),
                font_color=ICON_LABEL_COLOR, bold=False, alignment=PP_ALIGN.CENTER,
                font_name="Arial"):
    """Add a simple single-line text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None  # do not auto-resize
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    # Vertical alignment: top
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return txBox


# ---------------------------------------------------------------------------
# Add a family header bar
# ---------------------------------------------------------------------------
def add_family_header(slide, family_name, icon_count, y_pos):
    """Draw a coloured header bar with the family name and icon count."""
    left = Inches(PAGE_LEFT_MARGIN_IN)
    width = Inches(SLIDE_WIDTH_IN - 2 * PAGE_LEFT_MARGIN_IN)
    height = Inches(HEADER_HEIGHT_IN)

    # Background rectangle
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, y_pos, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = HEADER_BG_COLOR
    shape.line.fill.background()  # no border

    # Family name text
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    run = p.add_run()
    run.text = f"  {family_name}"
    run.font.size = HEADER_FONT_SIZE
    run.font.color.rgb = HEADER_TEXT_COLOR
    run.font.bold = True
    run.font.name = "Arial"

    # Icon count
    run2 = p.add_run()
    run2.text = f"   ({icon_count} icons)"
    run2.font.size = Pt(9)
    run2.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
    run2.font.bold = False
    run2.font.name = "Arial"

    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


# ---------------------------------------------------------------------------
# Layout calculation: how many icon rows fit on one slide?
# ---------------------------------------------------------------------------
def rows_per_slide(has_header):
    """Return how many icon rows fit on a slide, with or without a header."""
    available = USABLE_HEIGHT_IN
    if has_header:
        available -= HEADER_HEIGHT_IN + 0.05  # small gap after header
    return max(1, int(available / ROW_HEIGHT_IN))


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    print("=" * 60)
    print("OCHA Humanitarian Icons -- PowerPoint Generator")
    print("=" * 60)

    # --- Load metadata ---
    with open(METADATA_PATH, "r", encoding="utf-8") as f:
        metadata = json.load(f)

    families_order = metadata["families"]
    icons_data = metadata["icons"]

    # --- Group icons by family, sort alphabetically within each ---
    families = {}
    for icon_key, icon_info in icons_data.items():
        fam = icon_info.get("family", "Other")
        families.setdefault(fam, []).append((icon_key, icon_info))
    for fam in families:
        families[fam].sort(key=lambda x: x[1]["name"].lower())

    # --- Create presentation ---
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]  # blank

    total_icons = 0
    total_failed = 0
    total_slides = 0

    # --- Process each family ---
    for fam_name in families_order:
        if fam_name not in families:
            continue
        icon_list = families[fam_name]
        num_icons = len(icon_list)
        print(f"\n  Family: {fam_name} ({num_icons} icons)")

        idx = 0  # pointer into icon_list
        first_slide_of_family = True

        while idx < num_icons:
            slide = prs.slides.add_slide(blank_layout)
            total_slides += 1

            y_cursor = Inches(PAGE_TOP_MARGIN_IN)

            # Header on the first slide of each family
            if first_slide_of_family:
                add_family_header(slide, fam_name, num_icons, y_cursor)
                y_cursor += Inches(HEADER_HEIGHT_IN + 0.05)
                max_rows = rows_per_slide(has_header=True)
                first_slide_of_family = False
            else:
                # Continuation: add a lighter sub-header
                add_family_header(slide, f"{fam_name} (cont.)", num_icons, y_cursor)
                y_cursor += Inches(HEADER_HEIGHT_IN + 0.05)
                max_rows = rows_per_slide(has_header=True)

            # Fill rows
            rows_placed = 0
            while rows_placed < max_rows and idx < num_icons:
                # One row of icons
                for col in range(COLS):
                    if idx >= num_icons:
                        break

                    icon_key, icon_info = icon_list[idx]
                    idx += 1

                    # Cell position
                    cell_left = Inches(
                        PAGE_LEFT_MARGIN_IN + col * COL_WIDTH_IN
                    )
                    icon_top = y_cursor
                    label_top = icon_top + Inches(ICON_SIZE_IN + 0.02)

                    # Load and clean SVG
                    svg_path = os.path.join(SVG_DIR, icon_key + ".svg")
                    if not os.path.isfile(svg_path):
                        total_failed += 1
                        print(f"    WARNING: SVG not found: {icon_key}")
                        continue

                    result = clean_svg(svg_path, ICON_COLOR)
                    svg_doc, vb_w, vb_h = result
                    if svg_doc is None:
                        total_failed += 1
                        print(f"    WARNING: could not process: {icon_key}")
                        continue

                    svg_bytes = svg_doc.encode("utf-8")

                    # Calculate display size preserving aspect ratio
                    display_name = icon_info.get("name", icon_key)
                    if vb_w > 0 and vb_h > 0:
                        aspect = vb_w / vb_h
                        if aspect >= 1:
                            # Wider than tall: fit to width
                            disp_w = ICON_SIZE_IN
                            disp_h = ICON_SIZE_IN / aspect
                        else:
                            # Taller than wide: fit to height
                            disp_h = ICON_SIZE_IN
                            disp_w = ICON_SIZE_IN * aspect
                    else:
                        disp_w = ICON_SIZE_IN
                        disp_h = ICON_SIZE_IN

                    # Centre icon within cell
                    icon_left_adj = cell_left + Inches(
                        (COL_WIDTH_IN - disp_w) / 2
                    )
                    icon_top_adj = icon_top + Inches(
                        (ICON_SIZE_IN - disp_h) / 2
                    )

                    try:
                        add_svg_picture(
                            slide,
                            svg_bytes,
                            icon_left_adj,
                            icon_top_adj,
                            Inches(disp_w),
                            Inches(disp_h),
                            prs,
                            desc=display_name,
                        )
                    except Exception as exc:
                        total_failed += 1
                        print(f"    ERROR inserting {icon_key}: {exc}")
                        continue

                    # Label below icon
                    label_text = display_name
                    if len(label_text) > 20:
                        label_text = label_text[:19] + "\u2026"
                    add_textbox(
                        slide,
                        label_text,
                        cell_left,
                        label_top,
                        Inches(COL_WIDTH_IN),
                        Inches(0.30),
                        font_size=LABEL_FONT_SIZE,
                        font_color=ICON_LABEL_COLOR,
                        alignment=PP_ALIGN.CENTER,
                    )

                    total_icons += 1

                rows_placed += 1
                y_cursor += Inches(ROW_HEIGHT_IN)

    # --- Save ---
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    prs.save(OUTPUT_PATH)
    file_size = os.path.getsize(OUTPUT_PATH)

    # --- Report ---
    print("\n" + "=" * 60)
    print("RESULTS")
    print("=" * 60)
    print(f"  Output file : {OUTPUT_PATH}")
    print(f"  File size   : {file_size:,} bytes ({file_size / 1024:.1f} KB)")
    print(f"  Slides      : {total_slides}")
    print(f"  Icons OK    : {total_icons}")
    if total_failed:
        print(f"  Icons FAILED: {total_failed}")
    print(f"  Families    : {len([f for f in families_order if f in families])}")
    print(f"  Layout      : {SLIDE_WIDTH_IN}\" x {SLIDE_HEIGHT_IN}\" widescreen")
    print(f"  Grid        : {COLS} columns, ~{ICON_SIZE_IN}\" icons")
    print()
    print("  Icons are inserted as SVG images (with PNG fallback).")
    print("  To edit colors in PowerPoint:")
    print("    1. Right-click an icon")
    print("    2. Choose 'Convert to Shape'")
    print("    3. Then change fill colour as needed")
    print()


if __name__ == "__main__":
    main()
